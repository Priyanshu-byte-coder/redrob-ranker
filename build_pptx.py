"""
Build the hackathon submission PPTX from the Redrob template.
Preserves branding (background images, fonts) and fills content + diagrams.
Template has WHITE body background — all text must be DARK.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

TEMPLATE = Path(__file__).parent / "Idea Submission Template _ Redrob.pptx"
ASSETS = Path(__file__).parent / "assets"
OUTPUT = Path(__file__).parent / "Redrob_Submission_v5.pptx"

# Colors — DARK text on WHITE background
DARK = RGBColor(0x1E, 0x29, 0x3B)       # Primary text
MED = RGBColor(0x47, 0x55, 0x69)         # Secondary text
DIM = RGBColor(0x94, 0xA3, 0xB8)         # Muted/caption
ACCENT = RGBColor(0x4F, 0x46, 0xE5)      # Indigo headings
GREEN = RGBColor(0x05, 0x96, 0x69)
RED = RGBColor(0xDC, 0x26, 0x26)
# For slide 11 (dark background) — white text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCB, 0xD5, 0xE1)


def clear_textbox(shape):
    tf = shape.text_frame
    tf.clear()
    return tf


def add_p(tf, text, size=11, bold=False, color=DARK, space_after=Pt(3)):
    p = tf.add_paragraph()
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def first_p(tf, text, size=11, bold=False, color=DARK):
    p = tf.paragraphs[0]
    p.space_after = Pt(3)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_img(slide, path, left, top, width=None, height=None):
    kw = {"image_file": str(path), "left": left, "top": top}
    if width: kw["width"] = width
    if height: kw["height"] = height
    slide.shapes.add_picture(**kw)


def build():
    prs = Presentation(str(TEMPLATE))
    slides = list(prs.slides)

    # ===================== SLIDE 1: Title =====================
    s = slides[0]
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if "Team Name" in text:
            tf = clear_textbox(shape)
            first_p(tf, "Team Name :  Umbrella.co", size=16, bold=True, color=DARK)
        elif "Team Leader" in text:
            tf = clear_textbox(shape)
            first_p(tf, "Team Leader :  Priyanshu", size=16, bold=True, color=DARK)
        elif "Problem Statement" in text:
            tf = clear_textbox(shape)
            first_p(tf, "Track 1: Data & AI  |  Intelligent Candidate Discovery & Ranking", size=14, bold=True, color=DARK)

    # ===================== SLIDE 2: Solution Overview =====================
    s = slides[1]
    for shape in s.shapes:
        if shape.has_text_frame and "proposed solution" in shape.text_frame.text.lower():
            tf = clear_textbox(shape)
            tf.word_wrap = True

            first_p(tf, "What is your proposed solution?", size=12, bold=True, color=ACCENT)
            add_p(tf, "")
            add_p(tf, "A rule-based multi-signal ranking pipeline processing 100K candidates in ~25 seconds on CPU-only. Uses all 23/23 Redrob signals + career description analysis. No GPU, no LLM, no network — fully reproducible.", size=10, color=DARK)
            add_p(tf, "")
            add_p(tf, "4-Stage Pipeline:", size=10, bold=True, color=ACCENT)
            add_p(tf, "  Stage 0:  Honeypot Detection  \u2192  67 impossible profiles caught, 0 in top 100", size=9, color=MED)
            add_p(tf, "  Stage 1:  Coarse Filter  \u2192  100K \u2192 40,789 (title + description + experience)", size=9, color=MED)
            add_p(tf, "  Stage 2:  Multi-Signal Scoring  \u2192  8 weighted dimensions, 178-skill taxonomy", size=9, color=MED)
            add_p(tf, "  Stage 3:  Rank + Reasoning  \u2192  top 100 with fact-grounded per-candidate explanations", size=9, color=MED)
            add_p(tf, "")
            add_p(tf, "Why this beats embeddings:", size=12, bold=True, color=ACCENT)
            add_p(tf, "")
            add_p(tf, "Candidate A (embeddings rank HIGH \u2014 we rank LOW):", size=9, bold=True, color=RED)
            add_p(tf, "  \"Senior AI Engineer\" at TCS \u2022  Skills: PyTorch, BERT, RAG, FAISS, Weaviate", size=8, color=MED)
            add_p(tf, "  Descriptions: \u201cmanaged client deliverables\u201d, \u201cstakeholder management\u201d  \u2192  zero production ML", size=8, color=DIM)
            add_p(tf, "")
            add_p(tf, "Candidate B (embeddings rank LOW \u2014 we rank HIGH):", size=9, bold=True, color=GREEN)
            add_p(tf, "  \"Software Engineer\" at Razorpay \u2022  Skills: Python, PyTorch", size=8, color=MED)
            add_p(tf, "  Descriptions: \u201cbuilt recommendation pipeline 2M daily users\u201d, \u201cdeployed embedding model p99<50ms\u201d", size=8, color=DIM)

    # ===================== SLIDE 3: JD Understanding =====================
    s = slides[2]
    for shape in s.shapes:
        if shape.has_text_frame and "key requirements" in shape.text_frame.text.lower():
            tf = clear_textbox(shape)
            tf.word_wrap = True

            first_p(tf, "The JD told us how to solve it:", size=12, bold=True, color=ACCENT)
            add_p(tf, "")
            add_p(tf, "\u201cThe right answer is NOT find candidates whose skills section contains the most AI keywords.", size=9, bold=True, color=RED)
            add_p(tf, " That\u2019s a trap we\u2019ve explicitly built into the dataset.\u201d  \u2014 JD, Final Note", size=8, color=DIM)
            add_p(tf, "")
            add_p(tf, "What the JD actually requires:", size=11, bold=True, color=ACCENT)
            add_p(tf, "\u2022  Production ML deployments at product companies (not consulting)", size=9, color=DARK)
            add_p(tf, "\u2022  5\u20139yr exp; 0-to-1 builder mindset; leadership + startup background", size=9, color=DARK)
            add_p(tf, "\u2022  Embeddings, vector DBs, ranking systems \u2014 but only with production evidence", size=9, color=DARK)
            add_p(tf, "\u2022  Pune/Noida preferred; notice period \u226430 days ideal", size=9, color=DARK)
            add_p(tf, "\u2022  Behavioral availability: active on platform, responsive to recruiters", size=9, color=DARK)
            add_p(tf, "")
            add_p(tf, "How we read between the lines:", size=11, bold=True, color=ACCENT)
            add_p(tf, "1.  Career Trajectory \u2014 production ML evidence in descriptions (40+ keywords)", size=9, bold=True, color=DARK)
            add_p(tf, "2.  Trust Multiplier \u2014 Expert + 0mo duration = 0.3\u00d7 (kills keyword stuffers)", size=9, bold=True, color=DARK)
            add_p(tf, "3.  Anti-Patterns \u2014 consulting-only (\u22120.50), keyword stuffer (\u22120.80), title chaser (\u22120.20)", size=9, bold=True, color=DARK)
            add_p(tf, "4.  All 23 Signals \u2014 notice period, response rate, recency, salary feasibility, \u2026", size=9, bold=True, color=DARK)

    # ===================== SLIDE 4: Ranking Methodology =====================
    s = slides[3]
    for shape in s.shapes:
        if shape.has_text_frame and "retrieve" in shape.text_frame.text.lower():
            tf = clear_textbox(shape)
            tf.word_wrap = True

            first_p(tf, "How does the system retrieve, score, and rank?", size=12, bold=True, color=ACCENT)
            add_p(tf, "")
            add_p(tf, "Retrieve: Stream 100K JSONL \u2192 coarse filter to ~40K (title + description + exp)", size=9, color=DARK)
            add_p(tf, "Score: 8 weighted sub-scores computed per candidate \u2192 composite", size=9, color=DARK)
            add_p(tf, "Rank: Sort by composite, tie-break by candidate_id, generate reasoning", size=9, color=DARK)
            add_p(tf, "")
            add_p(tf, "Scoring Formula:", size=11, bold=True, color=ACCENT)
            add_p(tf, "Title(20%) + Skills(20%) + Career(20%) + Exp(10%)", size=9, bold=True, color=MED)
            add_p(tf, "+ Behavioral(10%) + Location(8%) + Anti-Pattern(7%) + Education(5%)", size=9, bold=True, color=MED)
            add_p(tf, "")
            add_p(tf, "Key algorithms and heuristics:", size=11, bold=True, color=ACCENT)
            add_p(tf, "\u2022  No ML models \u2014 pure rule-based, fully deterministic", size=9, color=DARK)
            add_p(tf, "\u2022  Bell curve for experience (centered at 7yr, range 5-9)", size=9, color=DARK)
            add_p(tf, "\u2022  3-tier skill taxonomy: 178 skills with trust multipliers", size=9, color=DARK)
            add_p(tf, "\u2022  Career description keyword analysis (40+ production ML terms)", size=9, color=DARK)
            add_p(tf, "\u2022  Company classification: 54 consulting + 80+ product + 20+ AI-native", size=9, color=DARK)
            add_p(tf, "\u2022  Leadership + startup bonuses for founding team fit", size=9, color=DARK)

    # Scoring chart — right side
    img = ASSETS / "scoring_weights.png"
    if img.exists():
        add_img(s, img, left=Emu(4800000), top=Emu(1350000), width=Emu(4100000))

    # ===================== SLIDE 5: Explainability =====================
    s = slides[4]
    for shape in s.shapes:
        if shape.has_text_frame and "ranking decisions" in shape.text_frame.text.lower():
            tf = clear_textbox(shape)
            tf.word_wrap = True

            first_p(tf, "How are ranking decisions explained?", size=12, bold=True, color=ACCENT)
            add_p(tf, "")
            add_p(tf, "Every candidate gets a specific, fact-based reasoning string:", size=9, color=DARK)
            add_p(tf, "")
            add_p(tf, "  Rank #1: \"Senior AI Engineer at Netflix with 7.8yr exp.", size=8, bold=True, color=MED)
            add_p(tf, "  3/3 roles in ML/AI; ML leadership; skills: LoRA, L2R, Weaviate, PEFT\"", size=8, color=DIM)
            add_p(tf, "")
            add_p(tf, "  Rank #5: \"Senior NLP Engineer at Mad Street Den with 8.0yr exp.", size=8, bold=True, color=MED)
            add_p(tf, "  Concerns: low recruiter response rate (16%)\"", size=8, color=DIM)
            add_p(tf, "")
            add_p(tf, "31/100 candidates include honest concern flags \u2014 we don't hide gaps.", size=9, bold=True, color=GREEN)
            add_p(tf, "")
            add_p(tf, "How do we prevent hallucinations?", size=12, bold=True, color=ACCENT)
            add_p(tf, "")
            add_p(tf, "\u2022  Zero LLM calls \u2014 every claim traces to a specific profile field", size=9, color=DARK)
            add_p(tf, "\u2022  Skills counted only if present in candidate's skills array", size=9, color=DARK)
            add_p(tf, "\u2022  ML job count from description keyword analysis, not inference", size=9, color=DARK)
            add_p(tf, "")
            add_p(tf, "How do we handle suspicious profiles?", size=12, bold=True, color=ACCENT)
            add_p(tf, "")
            add_p(tf, "\u2022  Honeypot detection: 6 rules \u2192 67 caught, 0 in top 100", size=9, color=DARK)
            add_p(tf, "\u2022  Trust multiplier: expert + 0mo = 0.3\u00d7 credit", size=9, color=DARK)
            add_p(tf, "\u2022  Keyword stuffer flag: AI skills + zero ML evidence in descriptions", size=9, color=DARK)
            add_p(tf, "\u2022  Date arithmetic catches fabricated career timelines", size=9, color=DARK)

    # ===================== SLIDE 6: Workflow (diagram) =====================
    s = slides[5]
    for shape in s.shapes:
        if shape.has_text_frame and "complete workflow" in shape.text_frame.text.lower():
            tf = clear_textbox(shape)
            first_p(tf, "", size=1)
    img = ASSETS / "pipeline_flow.png"
    if img.exists():
        add_img(s, img, left=Emu(350000), top=Emu(1050000), width=Emu(8450000))

    # ===================== SLIDE 7: Architecture (diagram) =====================
    s = slides[6]
    img = ASSETS / "architecture.png"
    if img.exists():
        add_img(s, img, left=Emu(350000), top=Emu(1050000), width=Emu(8450000))

    # ===================== SLIDE 8: Results (dashboard) =====================
    s = slides[7]
    for shape in s.shapes:
        if shape.has_text_frame and "results" in shape.text_frame.text.lower() and "insights" in shape.text_frame.text.lower():
            tf = clear_textbox(shape)
            first_p(tf, "", size=1)
    img = ASSETS / "results_dashboard.png"
    if img.exists():
        add_img(s, img, left=Emu(350000), top=Emu(1050000), width=Emu(8450000))

    # ===================== SLIDE 9: Technologies =====================
    s = slides[8]
    for shape in s.shapes:
        if shape.has_text_frame and "technologies" in shape.text_frame.text.lower() and "frameworks" in shape.text_frame.text.lower():
            tf = clear_textbox(shape)
            tf.word_wrap = True

            first_p(tf, "Core Stack:", size=12, bold=True, color=ACCENT)
            add_p(tf, "")
            add_p(tf, "\u2022  Python 3.11 \u2014 core language, zero ML library dependencies", size=10, color=DARK)
            add_p(tf, "\u2022  Rule-based scoring \u2014 1,900 lines of hand-crafted feature logic", size=10, color=DARK)
            add_p(tf, "\u2022  Streaming JSONL \u2014 memory-efficient, handles 100K in single pass", size=10, color=DARK)
            add_p(tf, "\u2022  Gradio 5.x \u2014 interactive sandbox on HuggingFace Spaces", size=10, color=DARK)
            add_p(tf, "\u2022  Git \u2014 22+ commits, iterative development with clear history", size=10, color=DARK)
            add_p(tf, "\u2022  Dependencies: only pyyaml + tqdm", size=10, color=DARK)
            add_p(tf, "")
            add_p(tf, "Why rule-based over ML/embeddings?", size=12, bold=True, color=ACCENT)
            add_p(tf, "")
            add_p(tf, "The JD explicitly warns: keyword matching is a trap.", size=10, bold=True, color=DARK)
            add_p(tf, "")
            add_p(tf, "\u2022  Embeddings = sophisticated keyword matching \u2014 can't read career context", size=9, color=MED)
            add_p(tf, "\u2022  LLMs = expensive, non-deterministic, hallucinate reasoning", size=9, color=MED)
            add_p(tf, "\u2022  Rules encode domain knowledge that embeddings can't learn from 1 JD", size=9, color=MED)
            add_p(tf, "\u2022  25 seconds on CPU, fully deterministic, every decision auditable", size=9, color=MED)

    # ===================== SLIDE 10: Submission Assets =====================
    s = slides[9]
    for shape in s.shapes:
        if shape.has_text_frame and "github" in shape.text_frame.text.lower():
            tf = clear_textbox(shape)
            tf.word_wrap = True

            first_p(tf, "GitHub Repository:", size=12, bold=True, color=ACCENT)
            add_p(tf, "github.com/Priyanshu-byte-coder/redrob-ranker", size=10, color=DARK)
            add_p(tf, "")
            add_p(tf, "Live Sandbox Demo:", size=12, bold=True, color=ACCENT)
            add_p(tf, "huggingface.co/spaces/bladebutcher/redrob-ranker", size=10, color=DARK)
            add_p(tf, "")
            add_p(tf, "Submission Output:", size=12, bold=True, color=ACCENT)
            add_p(tf, "submission.csv \u2014 100 candidates ranked with per-candidate reasoning", size=10, color=DARK)
            add_p(tf, "Score range: 0.776 \u2013 0.899  |  Top 10 avg: 0.872", size=9, color=MED)
            add_p(tf, "")
            add_p(tf, "Quality Assurance:", size=12, bold=True, color=ACCENT)
            add_p(tf, "\u2022  10 automated tests, all passing", size=10, color=DARK)
            add_p(tf, "\u2022  Submission format validator passes all checks", size=10, color=DARK)
            add_p(tf, "\u2022  0 honeypots in top 100 (67 detected total)", size=10, color=DARK)
            add_p(tf, "\u2022  Top 10 manually verified: all genuine AI/ML engineers at product companies", size=10, color=DARK)
            add_p(tf, "\u2022  25s runtime (limit: 5 min)  |  <2GB RAM (limit: 16GB)  |  CPU-only", size=10, color=DARK)

    # ===================== SLIDE 11: Thank You (DARK bg) =====================
    s = slides[10]
    txBox = s.shapes.add_textbox(Emu(1500000), Emu(2200000), Emu(6000000), Emu(2200000))
    tf = txBox.text_frame
    tf.word_wrap = True
    first_p(tf, "Umbrella.co", size=30, bold=True, color=WHITE)
    add_p(tf, "", size=8)
    add_p(tf, "Priyanshu  |  Solo Participant", size=16, bold=True, color=WHITE)
    add_p(tf, "", size=6)
    add_p(tf, "doshipriyanshu3@gmail.com", size=12, color=LIGHT)
    add_p(tf, "+91-9549926195", size=12, color=LIGHT)
    add_p(tf, "", size=6)
    add_p(tf, "India Runs Hackathon 2026  |  Track 1: Data & AI Challenge", size=11, color=LIGHT)

    # ===================== SAVE =====================
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build()
